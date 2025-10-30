import os
from flask import Flask, render_template, request, flash, redirect, url_for, send_file, Response
from werkzeug.utils import secure_filename
from docxtpl import DocxTemplate
from docx2pdf import convert
from PIL import Image, ImageDraw, ImageFont
import openpyxl
import pandas as pd
from datetime import datetime
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
from pptx.dml.color import RGBColor
import zipfile
import io
import mysql.connector

db_config = {
    'host': 'localhost',
    'user': 'root',
    'password': '',
    'database': 'user_result'
}

# Function to insert multiple certificates into MySQL
def insert_multiple_certificate(first_name, last_name, certificate_filename, certificate_file):
    conn = None
    cursor = None
    try:
        conn = mysql.connector.connect(**db_config)
        cursor = conn.cursor()
        query = "INSERT INTO certificates (first_name, last_name, certificate_filename, certificate_file) VALUES (%s, %s, %s, %s)"
        cursor.execute(query, (first_name, last_name, certificate_filename, certificate_file))
        conn.commit()
        print(f"✅ Inserted certificate for {first_name} {last_name}")
    except mysql.connector.Error as err:
        print(f"❌ Error inserting certificate into DB: {err}")
    finally:
        if cursor:
            cursor.close()
        if conn and conn.is_connected():
            conn.close()

# Function to insert individual certificate into MySQL
def insert_individual_certificate(first_name, last_name, certificate_filename, certificate_file):
    conn = None
    cursor = None
    try:
        conn = mysql.connector.connect(**db_config)
        cursor = conn.cursor()
        query = "INSERT INTO individual_certificate (first_name, last_name, certificate_filename, certificate_file) VALUES (%s, %s, %s, %s)"
        cursor.execute(query, (first_name, last_name, certificate_filename, certificate_file))
        conn.commit()
        print(f"✅ Inserted individual certificate for {first_name} {last_name}")
    except mysql.connector.Error as err:
        print(f"❌ Error inserting individual certificate into DB: {err}")
    finally:
        if cursor:
            cursor.close()
        if conn and conn.is_connected():
            conn.close()

# Function to insert a transcript into MySQL
def insert_transcript(name, filename, file_data):
    conn = None
    cursor = None
    try:
        conn = mysql.connector.connect(**db_config)
        cursor = conn.cursor()
        query = "INSERT INTO transcript (name, file_name, file) VALUES (%s, %s, %s)"
        cursor.execute(query, (name, filename, file_data))
        conn.commit()
        print(f"✅ Successfully inserted transcript for {name} into database")
        return True
    except mysql.connector.Error as err:
        print(f"❌ MySQL Error inserting transcript into DB: {err}")
        return False
    finally:
        if cursor:
            cursor.close()
        if conn and conn.is_connected():
            conn.close()

# Function to get all individual certificates from MySQL
def get_individual_certificates_from_db():
    conn = None
    cursor = None
    try:
        conn = mysql.connector.connect(**db_config)
        cursor = conn.cursor(dictionary=True)
        cursor.execute("SELECT * FROM individual_certificate")
        results = cursor.fetchall()
        return results
    except mysql.connector.Error as err:
        print(f"Error fetching individual certificates from DB: {err}")
        return []
    finally:
        if cursor:
            cursor.close()
        if conn and conn.is_connected():
            conn.close()

# Function to get all certificates from MySQL
def get_certificates_from_db():
    conn = None
    cursor = None
    try:
        conn = mysql.connector.connect(**db_config)
        cursor = conn.cursor(dictionary=True)
        cursor.execute("SELECT * FROM certificates")
        results = cursor.fetchall()
        return results
    except mysql.connector.Error as err:
        print(f"Error fetching certificates from DB: {err}")
        return []
    finally:
        if cursor:
            cursor.close()
        if conn and conn.is_connected():
            conn.close()

# Function to get all transcripts from MySQL
def get_transcripts_from_db():
    conn = None
    cursor = None
    try:
        conn = mysql.connector.connect(**db_config)
        cursor = conn.cursor(dictionary=True)
        cursor.execute("SELECT * FROM transcript")
        results = cursor.fetchall()
        return results
    except mysql.connector.Error as err:
        print(f"Error fetching transcripts from DB: {err}")
        return []
    finally:
        if cursor:
            cursor.close()
        if conn and conn.is_connected():
            conn.close()

# Function to verify database structure
def verify_database_structure():
    conn = None
    cursor = None
    try:
        conn = mysql.connector.connect(**db_config)
        cursor = conn.cursor(dictionary=True)
        
        # Check individual_certificate table
        cursor.execute("DESCRIBE individual_certificate")
        individual_columns = [col['Field'] for col in cursor.fetchall()]
        print(f"📋 individual_certificate columns: {individual_columns}")
        
        # Check certificates table
        cursor.execute("DESCRIBE certificates")
        certificates_columns = [col['Field'] for col in cursor.fetchall()]
        print(f"📋 certificates columns: {certificates_columns}")
        
        # Check transcript table
        cursor.execute("DESCRIBE transcript")
        transcript_columns = [col['Field'] for col in cursor.fetchall()]
        print(f"📋 transcript columns: {transcript_columns}")
        
    except mysql.connector.Error as err:
        print(f"❌ Database error: {err}")
    finally:
        if cursor:
            cursor.close()
        if conn and conn.is_connected():
            conn.close()

app = Flask(__name__)
app.secret_key = 'your_secret_key_here'
UPLOAD_FOLDER = 'uploads'
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

# Function for generating certificates
def generate_certificates(excel_file, template_file, output_folder, template_filename, option, font_path="arialbd.ttf", font_size=100):
    workbook = openpyxl.load_workbook(excel_file)
    sheet = workbook.active
    data = list(sheet.values)[1:]  # Skip header row
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)
    font_name = ImageFont.truetype(font_path, font_size)
    generated_files = []
    
    print(f"📊 Processing {len(data)} students for certificates...")
    
    for row in data:
        name = row[0]  # Assuming name is in the first column
        
        # Split name on space, with error handling
        name_parts = name.split(' ', 1)
        first_name = name_parts[0].upper() if name_parts else ""
        last_name = name_parts[1].upper() if len(name_parts) > 1 else ""
        
        output_path = None
        file_data = None
        
        print(f"🎓 Generating certificate for: {name}")
        
        if template_filename.lower().endswith('.png') or template_filename.lower().endswith('.jpg'):
            try:
                certificate = Image.open(template_file)
                width, height = certificate.size
                draw = ImageDraw.Draw(certificate)
                
                name_position = (width//2, (height//2)+40)
                draw.text(name_position, str(name), fill="navy", font=font_name, anchor='mm')
                output_path = os.path.join(output_folder, "certificate_" + str(name).replace(" ", "_") + ".png")
                certificate.save(output_path)
                
                # Read file content for database storage
                with open(output_path, 'rb') as f:
                    file_data = f.read()
                
                generated_files.append(output_path)
                print(f"✅ Certificate generated for {name}")
            except Exception as e:
                print(f"❌ Error generating image certificate for {name}: {e}")
        
        elif template_filename.lower().endswith(".pptx"):
            try:
                prs = Presentation(template_file)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text_frame"):
                            text_frame = shape.text_frame
                            for paragraph in text_frame.paragraphs:
                                full_text = paragraph.text
                                if "{{your_name}}" in full_text:
                                    paragraph.text = full_text.replace("{{your_name}}", str(name))
                                    for run in paragraph.runs:
                                        if str(name) in run.text:
                                            run.font.name = 'Arial'
                                            run.font.size = Pt(54)
                                            run.font.bold = True
                                            run.font.color.rgb = RGBColor(0, 0, 128)
                                    paragraph.alignment = PP_ALIGN.CENTER
                
                output_path_pptx = os.path.join(output_folder, "certificate_" + str(name).replace(" ", "_") + ".pptx")
                prs.save(output_path_pptx)
                output_path = output_path_pptx
                
                # Read file content for database storage
                with open(output_path_pptx, 'rb') as f:
                    file_data = f.read()
                
                generated_files.append(output_path)
                print(f"✅ PPTX certificate generated for {name}")
                
                if option in ['pdf', 'both']:
                    pdf_path = os.path.join(output_folder, "certificate_" + str(name).replace(" ", "_") + ".pdf")
                    convert(output_path_pptx, pdf_path)
                    
                    # Read PDF file content for database storage
                    with open(pdf_path, 'rb') as f:
                        pdf_file_data = f.read()
                    
                    if option == 'pdf':
                        os.remove(output_path_pptx)
                        output_path = pdf_path
                        file_data = pdf_file_data
                    generated_files.append(pdf_path)
                    print(f"✅ PDF certificate generated for {name}")
                    
            except Exception as e:
                print(f"❌ Error generating PPTX/PDF certificate for {name}: {e}")
        
        # Insert into DB for all template types
        if output_path and file_data:
            certificate_filename = os.path.basename(output_path)
            print(f"💾 Inserting into DB: first_name={first_name}, last_name={last_name}, filename={certificate_filename}")
            insert_multiple_certificate(first_name, last_name, certificate_filename, file_data)
    
    print(f"✅ All certificates have been generated! Total: {len(generated_files)} files")
    return generated_files

# Function for generating individual certificate
def generate_individual_certificate(name, template_file, output_folder, template_filename, font_path="arialbd.ttf", font_size=100):
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)
    font_name = ImageFont.truetype(font_path, font_size)
    file_data = None
    output_path = None
    
    if template_filename.lower().endswith('.png') or template_filename.lower().endswith('.jpg'):
        certificate = Image.open(template_file)
        width, height = certificate.size
        draw = ImageDraw.Draw(certificate)
        
        name_position = (width//2, (height//2)+40)
        draw.text(name_position, str(name), fill="navy", font=font_name, anchor='mm')
        output_path = os.path.join(output_folder, "certificate_" + str(name).replace(" ", "_") + ".png")
        certificate.save(output_path)
        
        # Read file content for database storage
        with open(output_path, 'rb') as f:
            file_data = f.read()
    
    return output_path, file_data

# Functions for generating Transcripts
def TranscriptExcel_data(filename):
    workbook = openpyxl.load_workbook(filename)
    sheet = workbook.active
    return list(sheet.values)

def TranscriptDocument(template, output_directory, row_data):
    doc = DocxTemplate(template)
    current_date = datetime.now().strftime("%B %d, %Y")
    doc.render({
        "student_id": row_data[0],
        "first_name": row_data[1],
        "last_name": row_data[2],
        "logic": row_data[3],
        "l_g": row_data[4],
        "bcum": row_data[5],
        "bc_g": row_data[6],
        "design": row_data[7],
        "d_g": row_data[8],
        "p1": row_data[9],
        "p1_g": row_data[10],
        "e1": row_data[11],
        "e1_g": row_data[12],
        "wd": row_data[13],
        "wd_g": row_data[14],
        "algo": row_data[15],
        "al_g": row_data[16],
        "p2": row_data[17],
        "p2_g": row_data[18],
        "e2": row_data[19],
        "e2_g": row_data[20],
        "sd": row_data[21],
        "sd_g": row_data[22],
        "js": row_data[23],
        "js_g": row_data[24],
        "php": row_data[25],
        "ph_g": row_data[26],
        "db": row_data[27],
        "db_g": row_data[28],
        "vc1": row_data[29],
        "v1_g": row_data[30],
        "node": row_data[31],
        "no_g": row_data[32],
        "e3": row_data[33],
        "e3_g": row_data[34],
        "p3": row_data[35],
        "p3_g": row_data[36],
        "oop": row_data[37],
        "op_g": row_data[38],
        "lar": row_data[39],
        "lar_g": row_data[40],
        "vue": row_data[41],
        "vu_g": row_data[42],
        "vc2": row_data[43],
        "v2_g": row_data[44],
        "e4": row_data[45],
        "e4_g": row_data[46],
        "p4": row_data[47],
        "p4_g": row_data[48],
        "int": row_data[49],
        "in_g": row_data[50],
        'cur_date': current_date
    })
    doc_name = os.path.join(output_directory, f"{row_data[1]}_{row_data[2]}.docx")
    doc.save(doc_name)
    return doc_name

def TranscriptPdf(doc_path, pdf_directory):
    pdf_path = os.path.join(pdf_directory, os.path.splitext(os.path.basename(doc_path))[0] + ".pdf")
    convert(doc_path, pdf_path)
    return pdf_path

def generate_transcripts(excel_file, template_file, docx_directory, pdf_directory, option):
    os.makedirs(docx_directory, exist_ok=True)
    os.makedirs(pdf_directory, exist_ok=True)
    data_rows = TranscriptExcel_data(excel_file)
    generated_files = []

    print(f"📊 Processing {len(data_rows[1:])} students for transcripts...")
    
    # Get current count of transcripts in database for verification
    initial_transcripts_count = len(get_transcripts_from_db())
    print(f"📊 Initial transcripts in database: {initial_transcripts_count}")

    for row in data_rows[1:]:
        student_name = f"{row[1]} {row[2]}"  # first_name + last_name
        print(f"🎓 Processing transcript for: {student_name}")
        
        if option in ["doc", "both"]:
            try:
                doc_path = TranscriptDocument(template_file, docx_directory, row)
                generated_files.append(doc_path)
                
                # Insert DOC transcript into database
                with open(doc_path, 'rb') as file:
                    file_data = file.read()
                
                filename = os.path.basename(doc_path)
                print(f"💾 Attempting to insert DOC transcript: {filename} for {student_name}")
                
                # Insert transcript with file content
                success = insert_transcript(student_name, filename, file_data)
                if success:
                    print(f"✅ DOC transcript inserted for {student_name}")
                else:
                    print(f"❌ Failed to insert DOC transcript for {student_name}")
                    
            except Exception as e:
                print(f"❌ Error generating DOC transcript for {student_name}: {e}")
        
        if option in ["pdf", "both"]:
            try:
                if option == "pdf":
                    # For PDF-only option, create doc temporarily then convert
                    doc_path = TranscriptDocument(template_file, pdf_directory, row)
                else:
                    # For "both" option, use the already created doc
                    doc_path = os.path.join(docx_directory, f"{row[1]}_{row[2]}.docx")
                    if not os.path.exists(doc_path):
                        doc_path = TranscriptDocument(template_file, docx_directory, row)
                
                pdf_path = TranscriptPdf(doc_path, pdf_directory)
                generated_files.append(pdf_path)
                
                # Insert PDF transcript into database
                with open(pdf_path, 'rb') as file:
                    file_data = file.read()
                
                filename = os.path.basename(pdf_path)
                print(f"💾 Attempting to insert PDF transcript: {filename} for {student_name}")
                
                # Insert transcript with file content
                success = insert_transcript(student_name, filename, file_data)
                if success:
                    print(f"✅ PDF transcript inserted for {student_name}")
                else:
                    print(f"❌ Failed to insert PDF transcript for {student_name}")
                
                # Clean up temporary doc file if we're only generating PDF
                if option == "pdf" and os.path.exists(doc_path):
                    os.remove(doc_path)
                    
            except Exception as e:
                print(f"❌ Error generating PDF transcript for {student_name}: {e}")
    
    # Final verification
    final_transcripts_count = len(get_transcripts_from_db())
    transcripts_added = final_transcripts_count - initial_transcripts_count
    print(f"📊 Final transcripts in database: {final_transcripts_count}")
    print(f"📊 Transcripts added this session: {transcripts_added}")
    
    print(f"✅ All files for option '{option}' have been generated! Total: {len(generated_files)} files")
    return generated_files

@app.route('/')
def home():
    return render_template('home.html')

@app.route('/complete-info-multiple')
def complete_info_multiple():
    return render_template('complete_certificate_info_multiple.html')

@app.route('/generate-multipule', methods=['POST'])
def generate_multiple():
    if request.method == 'POST':
        doc_type = request.form.get('doc_type')
        option = request.form.get('option')
        template = request.files.get('template')
        excel = request.files.get('excel')

        if not template or not excel:
            flash('Please upload both template and Excel files.')
            return redirect(url_for('complete_info_multiple'))
        
        template_filename = secure_filename(template.filename)
        excel_filename = secure_filename(excel.filename)
        template_path = os.path.join(app.config['UPLOAD_FOLDER'], template_filename)
        excel_path = os.path.join(app.config['UPLOAD_FOLDER'], excel_filename)
        template.save(template_path)
        excel.save(excel_path)
        
        try:
            if doc_type == 'transcript':
                docx_directory = 'Transcript_Doc'
                pdf_directory = 'Transcript_PDF'
                generated_files = generate_transcripts(excel_path, template_path, docx_directory, pdf_directory, option)
                
                # Verify transcripts were inserted into database
                transcripts_count = len(get_transcripts_from_db())
                flash(f'Transcripts generated successfully! {len(generated_files)} files created. {transcripts_count} transcripts in database.')
                
            elif doc_type == 'certificate':
                output_folder = 'Certificates'
                generated_files = generate_certificates(excel_path, template_path, output_folder, template_filename, option)
                
                # Verify certificates were inserted into database
                certificates_count = len(get_certificates_from_db())
                flash(f'Certificates generated successfully! {len(generated_files)} files created. {certificates_count} certificates in database.')
                
            else:
                flash('Invalid document type.')
                return redirect(url_for('complete_info_multiple'))
            
            return redirect(url_for('view'))
        except Exception as e:
            flash(f'An error occurred: {str(e)}')
            return redirect(url_for('complete_info_multiple'))

@app.route('/complete-info-certificate')
def complete_info_certificate():
    return render_template('complete-certificate.html')

@app.route('/complete-info-certificate-individual')
def complete_info_certificate_individual():
    return render_template('complete-certificate-individual.html')

@app.route('/generate_certificate_individual', methods=['POST'])
def generate_certificate_individual():
    first_name = request.form.get('first_name')
    last_name = request.form.get('last_name')
    
    if not first_name or not last_name:
        flash('Please provide both first and last name.')
        return redirect(url_for('complete_info_certificate_individual'))
    
    # Use default template
    template_path = os.path.join(app.root_path, 'template', 'certificate_template.png')
    template_filename = os.path.basename(template_path)
    
    if not os.path.exists(template_path):
        flash('Default certificate template not found. Please contact admin.')
        return redirect(url_for('complete_info_certificate_individual'))
    
    name = first_name.upper() + " " + last_name.upper()
    output_folder = 'Certificates_Individual'
    
    try:
        output_path, file_data = generate_individual_certificate(name, template_path, output_folder, template_filename)
        
        if output_path and file_data:
            certificate_filename = os.path.basename(output_path)
            
            # Insert into DB with file content
            insert_individual_certificate(first_name.upper(), last_name.upper(), certificate_filename, file_data)
            
            flash('Certificate generated successfully!')
        else:
            flash('Error generating certificate file.')
            
        return redirect(url_for('view'))
    except Exception as e:
        flash(f'Error generating certificate: {str(e)}')
        return redirect(url_for('complete_info_certificate_individual'))

@app.route('/complet-info-transcript')
def complete_info_transcript():
    return render_template('complete_transcript.html')

@app.route('/download-zip/<doc_type>')
def download_zip(doc_type):
    conn = None
    cursor = None
    
    try:
        conn = mysql.connector.connect(**db_config)
        cursor = conn.cursor(dictionary=True)
        
        zip_buffer = io.BytesIO()
        
        with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
            if doc_type == 'certificate':
                # Get all certificates from both tables
                cursor.execute("SELECT certificate_filename, certificate_file as file FROM certificates")
                certificates = cursor.fetchall()
                
                cursor.execute("SELECT certificate_filename, certificate_file as file FROM individual_certificate")
                individual_certificates = cursor.fetchall()
                
                all_certificates = certificates + individual_certificates
                
                for cert in all_certificates:
                    if cert['file']:
                        zip_file.writestr(cert['certificate_filename'], cert['file'])
                        
            elif doc_type == 'transcript':
                # Get all transcripts
                cursor.execute("SELECT file_name, file FROM transcript")
                transcripts = cursor.fetchall()
                
                for transcript in transcripts:
                    if transcript['file']:
                        zip_file.writestr(transcript['file_name'], transcript['file'])
            else:
                flash('Invalid type')
                return redirect(url_for('view'))
        
        zip_buffer.seek(0)
        return send_file(zip_buffer, as_attachment=True, download_name=f'{doc_type}_files.zip', mimetype='application/zip')
        
    except mysql.connector.Error as err:
        print(f"❌ Database error: {err}")
        flash('Database error occurred.')
        return redirect(url_for('view'))
    finally:
        if cursor:
            cursor.close()
        if conn and conn.is_connected():
            conn.close()

@app.route('/download_file/<dir_name>/<filename>')
def download_file(dir_name, filename):
    conn = None
    cursor = None
    
    try:
        conn = mysql.connector.connect(**db_config)
        cursor = conn.cursor(dictionary=True)
        
        # Determine which table to query based on dir_name
        if dir_name == 'Certificates_Individual':
            query = "SELECT certificate_file as file, certificate_filename FROM individual_certificate WHERE certificate_filename = %s"
        elif dir_name == 'Certificates':
            query = "SELECT certificate_file as file, certificate_filename FROM certificates WHERE certificate_filename = %s"
        elif dir_name == 'Transcript_PDF' or dir_name == 'Transcript_Doc':
            query = "SELECT file, file_name FROM transcript WHERE file_name = %s"
        else:
            flash('Invalid directory name.')
            return redirect(url_for('view'))
        
        cursor.execute(query, (filename,))
        result = cursor.fetchone()
        
        if result and result['file']:
            # Determine content type based on file extension
            content_type = 'application/octet-stream'
            if filename.lower().endswith('.pdf'):
                content_type = 'application/pdf'
            elif filename.lower().endswith(('.png', '.jpg', '.jpeg')):
                content_type = f'image/{filename.split(".")[-1].lower()}'
            elif filename.lower().endswith('.docx'):
                content_type = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
            elif filename.lower().endswith('.pptx'):
                content_type = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
            
            # Create a response with the file content from database
            response = Response(
                result['file'],
                mimetype=content_type,
                headers={
                    'Content-Disposition': f'attachment; filename={filename}'
                }
            )
            return response
        else:
            flash('File not found in database.')
            return redirect(url_for('view'))
            
    except mysql.connector.Error as err:
        print(f"❌ Database error: {err}")
        flash('Database error occurred.')
        return redirect(url_for('view'))
    finally:
        if cursor:
            cursor.close()
        if conn and conn.is_connected():
            conn.close()

# Debug route to check database content
@app.route('/debug-db')
def debug_db():
    """Debug route to check database content"""
    conn = None
    cursor = None
    try:
        conn = mysql.connector.connect(**db_config)
        cursor = conn.cursor(dictionary=True)
        
        # Check individual_certificate
        cursor.execute("SELECT nb, first_name, last_name, certificate_filename, LENGTH(certificate_file) as file_size FROM individual_certificate")
        individual_certs = cursor.fetchall()
        
        # Check certificates
        cursor.execute("SELECT nb, first_name, last_name, certificate_filename, LENGTH(certificate_file) as file_size FROM certificates")
        bulk_certs = cursor.fetchall()
        
        # Check transcript
        cursor.execute("SELECT nb, name, file_name, LENGTH(file) as file_size FROM transcript")
        transcripts = cursor.fetchall()
        
        debug_info = {
            'individual_certificates': individual_certs,
            'bulk_certificates': bulk_certs,
            'transcripts': transcripts
        }
        
        return f"""
        <h1>Database Debug Info</h1>
        <h2>Individual Certificates ({len(individual_certs)})</h2>
        <pre>{individual_certs}</pre>
        <h2>Bulk Certificates ({len(bulk_certs)})</h2>
        <pre>{bulk_certs}</pre>
        <h2>Transcripts ({len(transcripts)})</h2>
        <pre>{transcripts}</pre>
        """
        
    except mysql.connector.Error as err:
        return f"Database error: {err}"
    finally:
        if cursor:
            cursor.close()
        if conn and conn.is_connected():
            conn.close()

@app.route('/view')
def view():
    # Get data from all tables
    individual_certificates = get_individual_certificates_from_db()
    certificates = get_certificates_from_db()
    transcripts = get_transcripts_from_db()
    
    print(f"📊 Retrieved from DB - Individual Certificates: {len(individual_certificates)}, Certificates: {len(certificates)}, Transcripts: {len(transcripts)}")
    
    return render_template('view.html', 
                         individual_certificates=individual_certificates, 
                         certificates=certificates, 
                         transcripts=transcripts)

if __name__ == '__main__':
    # Verify database structure before starting
    verify_database_structure()
    app.run(debug=True)